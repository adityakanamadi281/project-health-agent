"""Builds the 5-7 slide monthly executive presentation from a MonthlySynthesis.

Uses python-pptx directly (no LLM involved here — by the time we reach this
module every fact has already been computed/synthesized). Colors and layout
are fixed constants for visual consistency, not stand-ins for data.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from project_health_agent.reporting.monthly_synthesis import MonthlySynthesis

NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)
RED = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xE6, 0x9A, 0x1D)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
GREY = RGBColor(0x6B, 0x6B, 0x6B)

RAG_COLOR = {"Green": GREEN, "Amber": AMBER, "Red": RED}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _text(slide, x, y, w, h, text, *, size=18, bold=False, color=CHARCOAL, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _bullets(slide, x, y, w, h, items, *, size=15, color=CHARCOAL, space_after=8):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(space_after)
    return box


def _footer(slide, text):
    _text(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.35), text, size=9, color=GREY)


def _title_slide(prs, synthesis: MonthlySynthesis):
    slide = _blank_slide(prs)
    _bg(slide, NAVY)
    _text(slide, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.2),
          "Portfolio Health Review", size=40, bold=True, color=WHITE)
    _text(slide, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.6),
          f"Professional Services — Monthly Synthesis · {synthesis.run_date.strftime('%B %Y')}",
          size=18, color=ICE)
    _text(slide, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.8),
          synthesis.headline, size=15, italic=True, color=ICE)
    return slide


def _rag_chart_slide(prs, synthesis: MonthlySynthesis):
    slide = _blank_slide(prs)
    _bg(slide, WHITE)
    _text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.7), "Portfolio at a Glance", size=28, bold=True, color=NAVY)

    dist = synthesis.stats.rag_distribution
    order = ["Green", "Amber", "Red"]
    chart_data = CategoryChartData()
    chart_data.categories = [c for c in order if c in dist] or order
    chart_data.add_series("Projects", [dist.get(c, 0) for c in (chart_data.categories)])

    x, y, w, h = Inches(0.7), Inches(1.4), Inches(6.3), Inches(5.3)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data)
    chart = gframe.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    for i, cat in enumerate(chart_data.categories):
        try:
            point = plot.series[0].points[i]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RAG_COLOR.get(cat, GREY)
        except IndexError:
            pass

    stat_x = Inches(7.4)
    _text(slide, stat_x, Inches(1.5), Inches(5.3), Inches(0.5), "Portfolio Snapshot", size=18, bold=True, color=NAVY)
    lines = [f"{synthesis.stats.project_count} active project(s) assessed this cycle"]
    if synthesis.stats.avg_schedule_variance_days is not None:
        lines.append(f"Average schedule variance: {synthesis.stats.avg_schedule_variance_days:.1f} days vs. baseline")
    if synthesis.stats.avg_sentiment_score is not None:
        lines.append(f"Average stakeholder sentiment: {synthesis.stats.avg_sentiment_score:+.2f} (-1 to +1)")
    if synthesis.stats.common_blocker_terms:
        top_terms = ", ".join(t for t, _ in synthesis.stats.common_blocker_terms[:3])
        lines.append(f"Recurring blocker themes: {top_terms}")
    if synthesis.stats.projects_with_perception_gap:
        lines.append("PM-vs-signal perception gap on: " + ", ".join(synthesis.stats.projects_with_perception_gap))
    _bullets(slide, stat_x, Inches(2.1), Inches(5.3), Inches(4.4), lines, size=14)
    _footer(slide, synthesis.stats.history_depth_note)
    return slide


def _trends_slide(prs, synthesis: MonthlySynthesis):
    slide = _blank_slide(prs)
    _bg(slide, WHITE)
    _text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.7), "Trends Across the Portfolio", size=28, bold=True, color=NAVY)
    _bullets(slide, Inches(0.7), Inches(1.4), Inches(11.8), Inches(5.3), synthesis.trend_bullets, size=17, space_after=14)
    return slide


def _risks_slide(prs, synthesis: MonthlySynthesis):
    slide = _blank_slide(prs)
    _bg(slide, WHITE)
    _text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.7), "Emerging Risks", size=28, bold=True, color=NAVY)
    _bullets(slide, Inches(0.7), Inches(1.4), Inches(11.8), Inches(5.3), synthesis.emerging_risks, size=17, color=RED, space_after=14)
    return slide


def _snapshot_table_slide(prs, synthesis: MonthlySynthesis):
    slide = _blank_slide(prs)
    _bg(slide, WHITE)
    _text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.7), "Project-by-Project Snapshot", size=28, bold=True, color=NAVY)

    rows = len(synthesis.project_snapshots) + 1
    cols = 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.5 * rows))
    table = table_shape.table
    headers = ["Project", "RAG", "Schedule Variance", "Summary"]
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(13)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE

    for r, snap in enumerate(synthesis.project_snapshots, start=1):
        table.cell(r, 0).text = snap.project_name
        table.cell(r, 1).text = snap.rag
        table.cell(r, 1).text_frame.paragraphs[0].font.color.rgb = RAG_COLOR.get(snap.rag, CHARCOAL)
        table.cell(r, 1).text_frame.paragraphs[0].font.bold = True
        variance_txt = f"{snap.schedule_variance_days:+d}d" if snap.schedule_variance_days is not None else "n/a"
        table.cell(r, 2).text = variance_txt
        table.cell(r, 3).text = snap.one_liner
        for c in range(cols):
            table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(12)
    return slide


def _recommendations_slide(prs, synthesis: MonthlySynthesis):
    slide = _blank_slide(prs)
    _bg(slide, NAVY)
    _text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.7), "Recommended Actions", size=28, bold=True, color=WHITE)
    _bullets(slide, Inches(0.7), Inches(1.5), Inches(11.8), Inches(5.0), synthesis.recommendations, size=18, color=ICE, space_after=16)
    _footer(slide, f"Narrative source: {synthesis.narrative_source}  ·  Generated by the Project Health Reporting Agent")
    return slide


def build_monthly_deck(synthesis: MonthlySynthesis, output_path: str | Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _title_slide(prs, synthesis)
    _rag_chart_slide(prs, synthesis)
    _trends_slide(prs, synthesis)
    _risks_slide(prs, synthesis)
    _snapshot_table_slide(prs, synthesis)
    _recommendations_slide(prs, synthesis)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
